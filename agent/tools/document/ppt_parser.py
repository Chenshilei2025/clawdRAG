"""
PPT content extraction tool.
"""
from typing import Dict, Any, List, Optional
from pathlib import Path
import logging

from ..base import Tool, ToolResult

logger = logging.getLogger(__name__)


class PPTContentExtractorTool(Tool):
    """Extract text, images, and structure from PowerPoint presentations."""

    name = "extract_ppt_content"
    description = "Extract content from PowerPoint files (.ppt, .pptx) including text, slide structure, and images."

    async def execute(
        self,
        file_path: str,
        include_images: bool = True,
        extract_notes: bool = True
    ) -> ToolResult:
        """Extract content from PPT file."""
        try:
            path = Path(file_path)
            if not path.exists():
                return ToolResult(success=False, error=f"File not found: {file_path}")

            if path.suffix.lower() not in {".ppt", ".pptx"}:
                return ToolResult(success=False, error=f"Not a PowerPoint file: {path.suffix}")

            # Use python-pptx for .pptx files
            if path.suffix.lower() == ".pptx":
                content = await self._extract_pptx(path, include_images, extract_notes)
            else:
                # For .ppt files, suggest conversion
                return ToolResult(
                    success=False,
                    error="Legacy .ppt format not supported. Please convert to .pptx first."
                )

            return ToolResult(
                success=True,
                data=content
            )

        except Exception as e:
            logger.error(f"Error extracting PPT content: {e}")
            return ToolResult(success=False, error=str(e))

    async def _extract_pptx(
        self,
        path: Path,
        include_images: bool,
        extract_notes: bool
    ) -> Dict[str, Any]:
        """Extract content from .pptx file."""
        try:
            from pptx import Presentation
        except ImportError:
            return {"error": "python-pptx not installed. Install with: pip install python-pptx"}

        prs = Presentation(str(path))

        slides_content = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_data = {
                "slide_number": slide_num,
                "title": "",
                "text_content": [],
                "notes": "",
                "images": []
            }

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        # Check if it's likely a title
                        if shape.shape_type == 14:  # Title shape type
                            slide_data["title"] = text
                        else:
                            slide_data["text_content"].append(text)

                # Extract images
                if include_images and shape.shape_type == 13:  # Picture shape type
                    slide_data["images"].append({
                        "type": "image",
                        "shape_id": shape.shape_id
                    })

            # Extract notes
            if extract_notes and slide.notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                slide_data["notes"] = notes_text

            slides_content.append(slide_data)

        return {
            "file_path": str(path),
            "total_slides": len(prs.slides),
            "slides": slides_content,
            "full_text": "\n\n".join([
                f"Slide {s['slide_number']}: {s['title']}\n" + "\n".join(s['text_content'])
                for s in slides_content
            ])
        }

    def get_parameters_schema(self) -> Dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Path to the PowerPoint file"
                },
                "include_images": {
                    "type": "boolean",
                    "description": "Include image information",
                    "default": True
                },
                "extract_notes": {
                    "type": "boolean",
                    "description": "Extract speaker notes",
                    "default": True
                }
            },
            "required": ["file_path"]
        }
